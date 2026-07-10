"""
DocSensei - Enterprise Multi-Document Adaptive RAG Assistant
==============================================================
A production-ready Flask application implementing an adaptive
Retrieval-Augmented Generation pipeline over user-uploaded documents
(PDF, DOCX, PPTX, TXT).

Pipeline: Dense Retrieval + BM25 -> Hybrid Fusion -> CrossEncoder Rerank
          -> Contextual Compression -> Groq Answer Generation

Run:
    python app.py

Required environment variables:
    PINECONE_API_KEY       - Pinecone API key
    PINECONE_INDEX_NAME    - Pinecone index name (default: "docsensei-index")
    PINECONE_CLOUD         - Pinecone serverless cloud (default: "aws")
    PINECONE_REGION        - Pinecone serverless region (default: "us-east-1")
    GROQ_API_KEY           - Groq API key
    PORT                   - Port to run on (default: 5000)

Optional environment variables:
    OCR_ENABLED            - "true"/"false", enable OCR fallback for scanned
                              PDFs (default: "true"). Requires the system
                              binaries `tesseract-ocr` and `poppler-utils`
                              to be installed on the host. If they are not
                              present, OCR is silently skipped and a warning
                              is logged -- normal (non-scanned) PDFs are
                              unaffected either way.

Deployment note:
    OCR requires system binaries (tesseract, poppler) that are available on
    Render and Railway (via apt buildpacks / Dockerfile) but are NOT
    installable on Vercel's serverless Python runtime. On Vercel, OCR will
    automatically and silently disable itself; everything else in this file
    (Flask, RAG pipeline, Pinecone, Groq) is unaffected and portable.

Author: DocSensei Engineering
"""

from __future__ import annotations

import io
import json
import logging
import os
import re
import shutil
import threading
import unicodedata
import uuid
from collections import deque
from datetime import datetime
from typing import Any

from flask import Flask, jsonify, render_template, request
from werkzeug.utils import secure_filename

# ---------------------------------------------------------------------------
# Third-party libraries
# ---------------------------------------------------------------------------
from pypdf import PdfReader
from docx import Document as DocxDocument
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.shapes.group import GroupShape
from sentence_transformers import CrossEncoder, SentenceTransformer
from rank_bm25 import BM25Okapi
from groq import Groq
from pinecone import Pinecone, ServerlessSpec
from dotenv import load_dotenv

load_dotenv()

# ---------------------------------------------------------------------------
# Logging configuration
# ---------------------------------------------------------------------------
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)-8s | %(name)s | %(message)s",
)
logger = logging.getLogger("DocSensei")

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOAD_FOLDER = os.path.join(BASE_DIR, "uploads")
TEMP_FOLDER = os.path.join(BASE_DIR, "tmp")
ALLOWED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt"}
MAX_UPLOAD_BYTES = 500 * 1024 * 1024  # 500 MB

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 150

EMBEDDING_MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"
EMBEDDING_DIMENSION = 384

RERANKER_MODEL_NAME = "cross-encoder/ms-marco-MiniLM-L-6-v2"
RERANK_CANDIDATE_POOL = 15   # how many hybrid results to feed the reranker
RERANK_TOP_N = 8             # how many reranked chunks survive to compression
COMPRESSION_TOP_N = 5        # only compress the top-5 reranked chunks

GROQ_MODEL = "llama-3.3-70b-versatile"
GROQ_TEMPERATURE = 0.1

PINECONE_API_KEY = os.environ.get("PINECONE_API_KEY", "")
PINECONE_INDEX_NAME = os.environ.get("PINECONE_INDEX_NAME", "docsensei-index")
PINECONE_CLOUD = os.environ.get("PINECONE_CLOUD", "aws")
PINECONE_REGION = os.environ.get("PINECONE_REGION", "us-east-1")

GROQ_API_KEY = os.environ.get("GROQ_API_KEY", "")

OCR_ENABLED = os.environ.get("OCR_ENABLED", "true").lower() == "true"
OCR_MIN_CHARS_PER_PAGE = 20  # below this, a PDF page is treated as "empty" -> OCR

MAX_CONVERSATION_TURNS = 6

QUERY_LABELS = (
    "Simple",
    "Complex",
    "Comparison",
    "Summarization",
    "Reasoning",
    "Multi-document",
    "Follow-up",
)

ROUTER_TOP_K = {
    "Simple": 3,
    "Complex": 8,
    "Comparison": 8,
    "Summarization": 10,
    "Reasoning": 8,
    "Multi-document": 10,
    "Follow-up": 5,
}

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(TEMP_FOLDER, exist_ok=True)

# ---------------------------------------------------------------------------
# Flask app
# ---------------------------------------------------------------------------
app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = MAX_UPLOAD_BYTES
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER

# ---------------------------------------------------------------------------
# Global clients (initialized lazily / at startup) -- singletons so repeated
# uploads/chats never reload expensive models or reconnect clients.
# ---------------------------------------------------------------------------
_embedding_model: SentenceTransformer | None = None
_reranker_model: CrossEncoder | None = None
_groq_client: Groq | None = None
_pinecone_client: Pinecone | None = None
_pinecone_index = None
_ocr_available: bool | None = None  # tri-state cache: None = not checked yet


def get_embedding_model() -> SentenceTransformer:
    """Lazily load the sentence-transformers embedding model (singleton)."""
    global _embedding_model
    if _embedding_model is None:
        logger.info("Loading embedding model: %s", EMBEDDING_MODEL_NAME)
        _embedding_model = SentenceTransformer(EMBEDDING_MODEL_NAME)
    return _embedding_model


def get_reranker_model() -> CrossEncoder:
    """Lazily load the CrossEncoder reranker model (singleton)."""
    global _reranker_model
    if _reranker_model is None:
        logger.info("Loading reranker model: %s", RERANKER_MODEL_NAME)
        _reranker_model = CrossEncoder(RERANKER_MODEL_NAME)
    return _reranker_model


def get_groq_client() -> Groq:
    """Lazily instantiate the Groq client (singleton)."""
    global _groq_client
    if _groq_client is None:
        if not GROQ_API_KEY:
            raise RuntimeError("GROQ_API_KEY environment variable is not set.")
        _groq_client = Groq(api_key=GROQ_API_KEY)
    return _groq_client


def get_pinecone_index():
    """Lazily initialize the Pinecone client and return the working index."""
    global _pinecone_client, _pinecone_index
    if _pinecone_index is not None:
        return _pinecone_index

    if not PINECONE_API_KEY:
        raise RuntimeError("PINECONE_API_KEY environment variable is not set.")

    logger.info("Initializing Pinecone client")
    _pinecone_client = Pinecone(api_key=PINECONE_API_KEY)

    existing_indexes = [idx["name"] for idx in _pinecone_client.list_indexes()]

    if PINECONE_INDEX_NAME not in existing_indexes:
        logger.info("Creating new Pinecone index: %s", PINECONE_INDEX_NAME)
        _pinecone_client.create_index(
            name=PINECONE_INDEX_NAME,
            dimension=EMBEDDING_DIMENSION,
            metric="cosine",
            spec=ServerlessSpec(cloud=PINECONE_CLOUD, region=PINECONE_REGION),
        )
    else:
        logger.info("Using existing Pinecone index: %s", PINECONE_INDEX_NAME)

    _pinecone_index = _pinecone_client.Index(PINECONE_INDEX_NAME)
    return _pinecone_index


def ocr_is_available() -> bool:
    """Check once (and cache) whether the OCR system binaries are usable."""
    global _ocr_available
    if _ocr_available is not None:
        return _ocr_available
    if not OCR_ENABLED:
        _ocr_available = False
        return False
    try:
        import pytesseract  # noqa: F401
        from pdf2image import convert_from_path  # noqa: F401
        # pytesseract.get_tesseract_version raises if the binary is missing.
        pytesseract.get_tesseract_version()
        _ocr_available = True
        logger.info("OCR fallback is available (tesseract + poppler detected)")
    except Exception as exc:
        _ocr_available = False
        logger.warning(
            "OCR fallback disabled (tesseract/poppler not found or pytesseract/"
            "pdf2image not installed): %s. Scanned/image-only PDFs will yield "
            "empty text for the affected pages.",
            exc,
        )
    return _ocr_available


# ---------------------------------------------------------------------------
# In-memory application state
# ---------------------------------------------------------------------------
class AppState:
    """Holds all mutable session-scoped state for the RAG pipeline.

    This is a single-tenant, in-memory store. It is intentionally simple
    (no external session backend) as specified for this deployment. Each
    upload is assigned a fresh Pinecone namespace so that vectors from
    different uploads never collide, and only the immediately-previous
    namespace is ever deleted.
    """

    def __init__(self) -> None:
        self.lock = threading.Lock()
        self.reset()

    def reset(self) -> None:
        with self.lock:
            self.uploaded_files: list[str] = []
            self.raw_documents: dict[str, str] = {}     # filename -> full extracted text
            self.chunks: list[dict[str, Any]] = []
            self.bm25_index: BM25Okapi | None = None
            self.bm25_corpus_tokens: list[list[str]] = []
            self.memory: deque = deque(maxlen=MAX_CONVERSATION_TURNS)
            self.summaries: dict[str, dict[str, Any]] = {}  # filename -> structured summary
            self.suggested_questions: dict[str, list[str]] = {}
            self.namespace: str | None = None            # current Pinecone namespace
            self.previous_namespace: str | None = None    # namespace to delete on next upload


state = AppState()

# ---------------------------------------------------------------------------
# SECTION: Document Loaders
# ---------------------------------------------------------------------------


def _ocr_pdf_page(path: str, page_number: int) -> str:
    """Render a single PDF page to an image and OCR it. Returns "" on failure."""
    try:
        import pytesseract
        from pdf2image import convert_from_path

        images = convert_from_path(
            path, first_page=page_number, last_page=page_number, dpi=200
        )
        if not images:
            return ""
        text = pytesseract.image_to_string(images[0])
        return text or ""
    except Exception:
        logger.exception("OCR failed for page %d of %s", page_number, path)
        return ""


def load_pdf(path: str) -> list[dict[str, Any]]:
    """Extract text from a PDF file, page by page, using pypdf.

    Pages with little or no extractable text (e.g. scanned images) are
    transparently retried through OCR when available, so the caller never
    needs to know whether a given page's text came from native extraction
    or OCR.
    """
    pages: list[dict[str, Any]] = []
    ocr_used_count = 0
    try:
        reader = PdfReader(path)
        for page_number, page in enumerate(reader.pages, start=1):
            try:
                text = page.extract_text() or ""
            except Exception:
                logger.exception(
                    "Native text extraction failed on page %d of %s; will try OCR",
                    page_number, path,
                )
                text = ""

            if len(text.strip()) < OCR_MIN_CHARS_PER_PAGE and ocr_is_available():
                ocr_text = _ocr_pdf_page(path, page_number)
                if len(ocr_text.strip()) > len(text.strip()):
                    text = ocr_text
                    ocr_used_count += 1

            if text.strip():
                pages.append({"text": text, "page": page_number})

        logger.info(
            "Loaded PDF '%s' with %d pages (OCR used on %d page(s))",
            os.path.basename(path), len(pages), ocr_used_count,
        )
    except Exception:
        logger.exception("Failed to load PDF: %s", path)
    return pages


def _iter_docx_headers_footers(doc: DocxDocument):
    """Yield (label, text) pairs for every header/footer across all sections."""
    for section_index, section in enumerate(doc.sections, start=1):
        for label, part in (
            (f"header (section {section_index})", section.header),
            (f"first-page header (section {section_index})", section.first_page_header),
            (f"footer (section {section_index})", section.footer),
            (f"first-page footer (section {section_index})", section.first_page_footer),
        ):
            try:
                text = "\n".join(p.text for p in part.paragraphs if p.text.strip())
            except Exception:
                text = ""
            if text.strip():
                yield label, text


def _extract_docx_textboxes(doc: DocxDocument) -> list[str]:
    """Extract text from DrawingML/VML text boxes via raw XML traversal.

    python-docx has no first-class API for text boxes, so we walk the
    underlying XML looking for <w:txbxContent> elements (used by both the
    legacy VML and modern DrawingML text box formats).
    """
    textbox_texts: list[str] = []
    try:
        body = doc.element.body
        for txbx in body.iter(qn("w:txbxContent")):
            paragraphs = txbx.iter(qn("w:p"))
            texts = []
            for p in paragraphs:
                runs = p.iter(qn("w:t"))
                run_text = "".join(t.text or "" for t in runs)
                if run_text.strip():
                    texts.append(run_text)
            joined = "\n".join(texts)
            if joined.strip():
                textbox_texts.append(joined)
    except Exception:
        logger.exception("Text box extraction failed for DOCX")
    return textbox_texts


def load_docx(path: str) -> list[dict[str, Any]]:
    """Extract text from a DOCX file: paragraphs, headings, tables, headers,
    footers, and (best-effort) text boxes. Empty objects are skipped.
    """
    pages: list[dict[str, Any]] = []
    try:
        doc = DocxDocument(path)
        parts: list[str] = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style_name = (para.style.name or "").lower() if para.style else ""
            if "heading" in style_name or "title" in style_name:
                parts.append(f"## {text}")
            else:
                parts.append(text)

        for table_index, table in enumerate(doc.tables, start=1):
            table_rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    table_rows.append(" | ".join(cells))
            if table_rows:
                parts.append(f"[Table {table_index}]\n" + "\n".join(table_rows))

        for label, text in _iter_docx_headers_footers(doc):
            parts.append(f"[{label}] {text}")

        for i, textbox_text in enumerate(_extract_docx_textboxes(doc), start=1):
            parts.append(f"[Text box {i}] {textbox_text}")

        full_text = "\n".join(p for p in parts if p.strip())
        if full_text.strip():
            pages.append({"text": full_text, "page": 1})
        logger.info("Loaded DOCX '%s' (%d content blocks)", os.path.basename(path), len(parts))
    except Exception:
        logger.exception("Failed to load DOCX: %s", path)
    return pages


def _iter_pptx_shapes(shapes):
    """Recursively yield shapes, descending into grouped shapes."""
    for shape in shapes:
        if isinstance(shape, GroupShape) or getattr(shape, "shape_type", None) == 6:
            try:
                yield from _iter_pptx_shapes(shape.shapes)
                continue
            except Exception:
                pass
        yield shape


def _extract_pptx_chart_labels(shape) -> list[str]:
    """Extract category/series labels from a chart shape, if present."""
    labels: list[str] = []
    try:
        if not shape.has_chart:
            return labels
        chart = shape.chart
        for plot in chart.plots:
            try:
                categories = [str(c) for c in plot.categories if c is not None]
                labels.extend(categories)
            except Exception:
                pass
            for series in plot.series:
                name = getattr(series, "name", None)
                if name:
                    labels.append(str(name))
    except Exception:
        logger.exception("Chart label extraction failed")
    return labels


def load_pptx(path: str) -> list[dict[str, Any]]:
    """Extract text from a PPTX file per slide: titles, placeholders, bullet
    lists, grouped shapes, text boxes, speaker notes, and chart labels.
    Purely decorative objects (no text/data) are ignored.
    """
    pages: list[dict[str, Any]] = []
    try:
        presentation = Presentation(path)
        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_parts: list[str] = []

            for shape in _iter_pptx_shapes(slide.shapes):
                if getattr(shape, "has_text_frame", False):
                    is_title = False
                    try:
                        is_title = shape.placeholder_format is not None and shape.placeholder_format.idx == 0
                    except Exception:
                        is_title = False
                    for paragraph in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in paragraph.runs).strip()
                        if not line and paragraph.text.strip():
                            line = paragraph.text.strip()
                        if line:
                            slide_parts.append(f"# {line}" if is_title else line)

                chart_labels = _extract_pptx_chart_labels(shape)
                if chart_labels:
                    slide_parts.append("[Chart] " + ", ".join(chart_labels))

            try:
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_parts.append(f"[Speaker notes] {notes_text}")
            except Exception:
                logger.exception("Speaker notes extraction failed on slide %d", slide_number)

            slide_text = "\n".join(slide_parts)
            if slide_text.strip():
                pages.append({"text": slide_text, "page": slide_number})

        logger.info("Loaded PPTX '%s' with %d slides of text", os.path.basename(path), len(pages))
    except Exception:
        logger.exception("Failed to load PPTX: %s", path)
    return pages


def load_txt(path: str) -> list[dict[str, Any]]:
    """Load a plain text file."""
    pages: list[dict[str, Any]] = []
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        if text.strip():
            pages.append({"text": text, "page": 1})
        logger.info("Loaded TXT '%s'", os.path.basename(path))
    except Exception:
        logger.exception("Failed to load TXT: %s", path)
    return pages


LOADER_MAP = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".txt": load_txt,
}

# ---------------------------------------------------------------------------
# SECTION: Preprocessing
# ---------------------------------------------------------------------------


def clean_text(text: str) -> str:
    """Normalize unicode, collapse whitespace, and remove duplicate blank
    lines and empty paragraphs, while preserving useful structural markers
    (headings, bullets) produced by the loaders above.
    """
    if not text:
        return ""
    text = unicodedata.normalize("NFKC", text)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    lines = [line.strip() for line in text.split("\n")]
    deduped_lines = []
    previous_blank = False
    for line in lines:
        if line == "":
            if previous_blank:
                continue
            previous_blank = True
        else:
            previous_blank = False
        deduped_lines.append(line)
    return "\n".join(deduped_lines).strip()


# ---------------------------------------------------------------------------
# SECTION: Chunking (custom recursive character splitter)
# ---------------------------------------------------------------------------

BULLET_PREFIXES = ("- ", "* ", "\u2022 ", "\u25cf ")


def _is_heading_line(line: str) -> bool:
    return line.startswith("#") or line.startswith("[Table") or line.startswith("[Speaker notes")


def _is_bullet_line(line: str) -> bool:
    return line.lstrip().startswith(BULLET_PREFIXES)


def _presegment_structure_aware(text: str) -> list[str]:
    """Group lines into structural blocks so that headings stay attached to
    the paragraph that follows them, and consecutive bullet-list lines stay
    together as one block, before the recursive splitter ever sees them.
    """
    lines = text.split("\n")
    blocks: list[str] = []
    current: list[str] = []
    current_is_bullets = False

    def flush():
        if current:
            blocks.append("\n".join(current))
            current.clear()

    for line in lines:
        if not line.strip():
            flush()
            current_is_bullets = False
            continue
        if _is_heading_line(line):
            flush()
            current.append(line)
            current_is_bullets = False
            continue
        if _is_bullet_line(line):
            if current and not current_is_bullets and not _is_heading_line(current[-1]):
                flush()
            current.append(line)
            current_is_bullets = True
            continue
        # Regular paragraph line
        if current_is_bullets:
            flush()
        current.append(line)
        current_is_bullets = False

    flush()
    return [b for b in blocks if b.strip()]


def _split_by_separator(text: str, separator: str) -> list[str]:
    if separator == "":
        return list(text)
    return text.split(separator)


def recursive_character_split(
    text: str, chunk_size: int = CHUNK_SIZE, chunk_overlap: int = CHUNK_OVERLAP
) -> list[str]:
    """A dependency-free recursive character splitter modeled after
    LangChain's RecursiveCharacterTextSplitter, using a cascading list
    of separators from coarse to fine granularity. Structural blocks
    (headings + following paragraph, bullet-list runs) are pre-merged so
    the splitter is much less likely to sever a heading from its body
    text or break a bullet list mid-list.
    """
    separators = ["\n\n", "\n", ". ", " ", ""]

    def _split(text_block: str, seps: list[str]) -> list[str]:
        if len(text_block) <= chunk_size:
            return [text_block] if text_block.strip() else []

        if not seps:
            return [text_block[i:i + chunk_size] for i in range(0, len(text_block), chunk_size)]

        sep = seps[0]
        pieces = _split_by_separator(text_block, sep)

        chunks: list[str] = []
        current = ""
        for piece in pieces:
            candidate = current + (sep if current else "") + piece if sep else current + piece
            if len(candidate) <= chunk_size:
                current = candidate
            else:
                if current.strip():
                    chunks.append(current)
                if len(piece) > chunk_size:
                    chunks.extend(_split(piece, seps[1:]))
                    current = ""
                else:
                    current = piece
        if current.strip():
            chunks.append(current)
        return chunks

    # Pre-segment into structural blocks, then merge small adjacent blocks
    # up to chunk_size before falling back to the generic recursive split.
    blocks = _presegment_structure_aware(text)
    merged_blocks: list[str] = []
    current = ""
    for block in blocks:
        candidate = (current + "\n\n" + block) if current else block
        if len(candidate) <= chunk_size:
            current = candidate
        else:
            if current:
                merged_blocks.append(current)
            current = block
    if current:
        merged_blocks.append(current)

    raw_chunks: list[str] = []
    for block in merged_blocks:
        raw_chunks.extend(_split(block, separators))

    # Apply overlap
    overlapped_chunks = []
    for i, chunk in enumerate(raw_chunks):
        if i == 0 or chunk_overlap <= 0:
            overlapped_chunks.append(chunk)
        else:
            prev_tail = raw_chunks[i - 1][-chunk_overlap:]
            overlapped_chunks.append((prev_tail + chunk)[:chunk_size + chunk_overlap])
    return [c.strip() for c in overlapped_chunks if c.strip()]


def build_chunks_for_file(filename: str, file_type: str, pages: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Convert loaded page/slide text blocks into metadata-rich chunks."""
    chunks = []
    for page_entry in pages:
        cleaned = clean_text(page_entry["text"])
        if not cleaned:
            continue
        text_pieces = recursive_character_split(cleaned)
        for piece in text_pieces:
            chunk_id = str(uuid.uuid4())
            chunks.append({
                "chunk_id": chunk_id,
                "text": piece,
                "source_file": filename,
                "page": page_entry["page"],
                "file_type": file_type,
            })
    logger.info("Built %d chunks for file '%s'", len(chunks), filename)
    return chunks


# ---------------------------------------------------------------------------
# SECTION: Vector Store (Pinecone) Operations -- namespace-per-upload
# ---------------------------------------------------------------------------


def new_namespace() -> str:
    """Generate a fresh, collision-resistant namespace for one upload batch."""
    return f"ns-{uuid.uuid4().hex}"


def delete_namespace() -> None:
    """Delete all vectors from the Pinecone index."""
    try:
        index = get_pinecone_index()
        index.delete(delete_all=True)
        logger.info("Deleted all vectors from Pinecone index")
    except Exception as exc:
        logger.warning("Pinecone delete warning: %s", exc)


def embed_texts(texts: list[str]) -> list[list[float]]:
    """Embed a list of texts using the sentence-transformers model, batched
    and L2-normalized so cosine similarity behaves as expected in Pinecone.
    """
    if not texts:
        return []
    model = get_embedding_model()
    embeddings = model.encode(
        texts,
        batch_size=32,
        normalize_embeddings=True,
        show_progress_bar=False,
    )
    return embeddings.tolist()


def upsert_chunks_to_pinecone(chunks: list[dict[str, Any]], namespace: str) -> None:
    """Embed and upsert chunks into Pinecone in batches, scoped to namespace."""
    if not chunks:
        return
    index = get_pinecone_index()
    texts = [c["text"] for c in chunks]
    embeddings = embed_texts(texts)

    vectors = []
    for chunk, vector in zip(chunks, embeddings):
        vectors.append({
            "id": chunk["chunk_id"],
            "values": vector,
            "metadata": {
                "text": chunk["text"][:4000],  # metadata size guard
                "source_file": chunk["source_file"],
                "page": chunk["page"],
                "file_type": chunk["file_type"],
                "chunk_id": chunk["chunk_id"],
            },
        })

    batch_size = 100
    for i in range(0, len(vectors), batch_size):
        batch = vectors[i:i + batch_size]
        index.upsert(vectors=batch, namespace=namespace)
    logger.info("Upserted %d vectors into Pinecone namespace '%s'", len(vectors), namespace)


def dense_retrieve(query: str, top_k: int = 5) -> list[dict[str, Any]]:
    """Retrieve the top_k most relevant chunks via dense vector search,
    scoped to the current upload's namespace.
    """
    if not state.namespace:
        return []
    index = get_pinecone_index()
    query_vector = embed_texts([query])[0]
    result = index.query(
        vector=query_vector, top_k=top_k, include_metadata=True, namespace=state.namespace
    )
    matches = result.get("matches", []) if isinstance(result, dict) else result.matches
    retrieved = []
    for match in matches:
        metadata = match["metadata"] if isinstance(match, dict) else match.metadata
        score = match["score"] if isinstance(match, dict) else match.score
        retrieved.append({
            "chunk_id": metadata.get("chunk_id"),
            "text": metadata.get("text", ""),
            "source_file": metadata.get("source_file"),
            "page": metadata.get("page"),
            "file_type": metadata.get("file_type"),
            "score": float(score),
            "retriever": "dense",
        })
    return retrieved


# ---------------------------------------------------------------------------
# SECTION: BM25 Retrieval
# ---------------------------------------------------------------------------


def _tokenize(text: str) -> list[str]:
    return re.findall(r"[a-z0-9]+", text.lower())


def build_bm25_index() -> None:
    """Build (or rebuild) the in-memory BM25 index from the current chunks.
    Called automatically after every upload so BM25 always reflects the
    latest document set.
    """
    if not state.chunks:
        state.bm25_index = None
        state.bm25_corpus_tokens = []
        return
    corpus_tokens = [_tokenize(chunk["text"]) for chunk in state.chunks]
    state.bm25_corpus_tokens = corpus_tokens
    state.bm25_index = BM25Okapi(corpus_tokens)
    logger.info("Built BM25 index over %d chunks", len(state.chunks))


def bm25_retrieve(query: str, top_k: int = 5) -> list[dict[str, Any]]:
    """Retrieve top_k chunks via BM25 sparse keyword search."""
    if state.bm25_index is None or not state.chunks:
        return []
    query_tokens = _tokenize(query)
    scores = state.bm25_index.get_scores(query_tokens)
    ranked_indices = sorted(range(len(scores)), key=lambda i: scores[i], reverse=True)[:top_k]

    retrieved = []
    for idx in ranked_indices:
        if scores[idx] <= 0:
            continue
        chunk = state.chunks[idx]
        retrieved.append({
            "chunk_id": chunk["chunk_id"],
            "text": chunk["text"],
            "source_file": chunk["source_file"],
            "page": chunk["page"],
            "file_type": chunk["file_type"],
            "score": float(scores[idx]),
            "retriever": "bm25",
        })
    return retrieved


# ---------------------------------------------------------------------------
# SECTION: Hybrid / Ensemble Retrieval
# ---------------------------------------------------------------------------


def _min_max_normalize(results: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Min-max normalize scores within a result set to [0, 1]. Falls back to
    a constant 1.0 when all scores are equal, avoiding divide-by-zero and
    avoiding the previous max-only normalization's bias toward the top hit.
    """
    if not results:
        return results
    scores = [r["score"] for r in results]
    lo, hi = min(scores), max(scores)
    spread = hi - lo
    for r in results:
        r["norm_score"] = 1.0 if spread == 0 else (r["score"] - lo) / spread
    return results


def hybrid_retrieve(
    query: str, top_k: int = 5, dense_weight: float = 0.6, bm25_weight: float = 0.4
) -> list[dict[str, Any]]:
    """Combine dense and BM25 retrieval into a single ranked, deduplicated
    list using min-max normalized, weighted score fusion.
    """
    dense_results = dense_retrieve(query, top_k=top_k * 2)
    bm25_results = bm25_retrieve(query, top_k=top_k * 2)

    dense_results = _min_max_normalize(dense_results)
    bm25_results = _min_max_normalize(bm25_results)

    combined: dict[str, dict[str, Any]] = {}
    for r in dense_results:
        combined[r["chunk_id"]] = dict(r, weighted_score=r["norm_score"] * dense_weight)
    for r in bm25_results:
        if r["chunk_id"] in combined:
            combined[r["chunk_id"]]["weighted_score"] += r["norm_score"] * bm25_weight
            combined[r["chunk_id"]]["retriever"] = "hybrid"
        else:
            combined[r["chunk_id"]] = dict(r, weighted_score=r["norm_score"] * bm25_weight)

    ranked = sorted(combined.values(), key=lambda r: r["weighted_score"], reverse=True)
    return ranked[:top_k]


# ---------------------------------------------------------------------------
# SECTION: CrossEncoder Reranking
# ---------------------------------------------------------------------------


def rerank_chunks(query: str, chunks: list[dict[str, Any]], top_n: int = RERANK_TOP_N) -> list[dict[str, Any]]:
    """Rerank hybrid-retrieved chunks with a CrossEncoder for much higher
    precision than embedding similarity alone. Falls back to the input
    order (already hybrid-ranked) if the reranker fails to load or run.
    """
    if not chunks:
        return chunks
    try:
        model = get_reranker_model()
        pairs = [[query, c["text"]] for c in chunks]
        scores = model.predict(pairs)
        for chunk, score in zip(chunks, scores):
            chunk["rerank_score"] = float(score)
        ranked = sorted(chunks, key=lambda c: c["rerank_score"], reverse=True)
        return ranked[:top_n]
    except Exception:
        logger.exception("Reranking failed; falling back to hybrid order")
        return chunks[:top_n]


# ---------------------------------------------------------------------------
# SECTION: Adaptive Router (LLM-based)
# ---------------------------------------------------------------------------

ROUTER_SYSTEM_PROMPT = (
    "You are a query classification router for a document question-answering "
    "system. Classify the user's query into EXACTLY ONE of these labels:\n"
    "- Simple: a direct factual lookup answerable from one short passage.\n"
    "- Complex: requires synthesizing multiple ideas or reasoning steps.\n"
    "- Comparison: asks to compare, contrast, or weigh two or more things.\n"
    "- Summarization: asks for an overview, summary, or high-level recap.\n"
    "- Reasoning: asks 'why' or requires inference beyond stated facts.\n"
    "- Multi-document: explicitly or implicitly spans multiple documents.\n"
    "- Follow-up: only makes sense in light of prior conversation turns.\n"
    "Respond with ONLY the single label text, nothing else."
)


def _heuristic_classify_query(query: str) -> str:
    """Fast, dependency-free fallback classifier used only if the LLM
    router call fails (e.g. Groq outage), so retrieval never hard-stops.
    """
    q = query.strip().lower()
    word_count = len(q.split())
    if q.split() and q.split()[0] in ("it", "that", "this", "they", "them", "those"):
        return "Follow-up"
    if any(m in q for m in ("compare", "difference", "versus", " vs ", "contrast")):
        return "Comparison"
    if any(m in q for m in ("summarize", "summary", "overview", "recap")):
        return "Summarization"
    if any(m in q for m in ("why", "explain", "reasoning", "how come")):
        return "Reasoning"
    if word_count > 20 or "across" in q or "all documents" in q:
        return "Multi-document"
    if word_count <= 3:
        return "Follow-up"
    return "Simple"


def classify_query(query: str) -> str:
    """Classify a query into one of QUERY_LABELS using an LLM router, with
    a heuristic fallback if the Groq call fails for any reason.
    """
    try:
        client = get_groq_client()
        response = client.chat.completions.create(
            model=GROQ_MODEL,
            temperature=0.0,
            messages=[
                {"role": "system", "content": ROUTER_SYSTEM_PROMPT},
                {"role": "user", "content": f"Query: {query}\n\nLabel:"},
            ],
        )
        label = response.choices[0].message.content.strip()
        # Normalize against known labels defensively (LLM may add stray punctuation).
        for known in QUERY_LABELS:
            if known.lower() == label.lower().strip(". "):
                return known
        logger.warning("Router returned unrecognized label '%s'; using heuristic fallback", label)
        return _heuristic_classify_query(query)
    except Exception:
        logger.exception("LLM router call failed; using heuristic fallback")
        return _heuristic_classify_query(query)


def adaptive_retrieve(query: str) -> tuple[list[dict[str, Any]], str]:
    """Route the query through the adaptive LLM router, retrieve via hybrid
    search, then rerank with the CrossEncoder for final precision ordering.
    """
    query_type = classify_query(query)
    top_k = ROUTER_TOP_K.get(query_type, 5)
    logger.info("Adaptive router classified query as '%s' (top_k=%d)", query_type, top_k)

    hybrid_results = hybrid_retrieve(query, top_k=max(top_k, RERANK_CANDIDATE_POOL))
    reranked = rerank_chunks(query, hybrid_results, top_n=min(top_k, RERANK_TOP_N))
    return reranked, query_type


# ---------------------------------------------------------------------------
# SECTION: Contextual Compression
# ---------------------------------------------------------------------------


def compress_context(query: str, retrieved_chunks: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Use the LLM to extract only the sentences relevant to the query from
    each of the top-N reranked chunks, discarding irrelevant filler content.
    Only the top COMPRESSION_TOP_N chunks are compressed; any remaining
    lower-ranked chunks are passed through unmodified to avoid wasting LLM
    calls on context that contributes little to the final answer.
    """
    if not retrieved_chunks:
        return retrieved_chunks

    to_compress = retrieved_chunks[:COMPRESSION_TOP_N]
    passthrough = retrieved_chunks[COMPRESSION_TOP_N:]

    client = get_groq_client()
    numbered_chunks = "\n\n".join(
        f"[CHUNK {i}]\n{c['text']}" for i, c in enumerate(to_compress)
    )

    system_prompt = (
        "You are a context compression engine. For each numbered CHUNK, extract ONLY "
        "the sentences or phrases directly relevant to the user's question. If a chunk "
        "has nothing relevant, output an empty string for it. "
        "Respond ONLY with a JSON object mapping chunk index (as string) to the extracted text, "
        "with no preamble, no markdown fences, and no extra commentary."
    )
    user_prompt = f"Question: {query}\n\n{numbered_chunks}"

    try:
        response = client.chat.completions.create(
            model=GROQ_MODEL,
            temperature=GROQ_TEMPERATURE,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
        )
        raw = response.choices[0].message.content.strip()
        raw = re.sub(r"^```(json)?|```$", "", raw, flags=re.MULTILINE).strip()
        extracted_map = json.loads(raw)
    except Exception:
        logger.exception("Contextual compression failed; falling back to raw chunks")
        return retrieved_chunks

    compressed = []
    for i, chunk in enumerate(to_compress):
        extracted = extracted_map.get(str(i), "").strip()
        new_chunk = dict(chunk)
        new_chunk["text"] = extracted if extracted else chunk["text"]
        compressed.append(new_chunk)
    return compressed + passthrough


# ---------------------------------------------------------------------------
# SECTION: Query Rewriting
# ---------------------------------------------------------------------------

STANDALONE_CHECK_PROMPT = (
    "Determine whether the following follow-up question is already a fully "
    "standalone question that can be understood WITHOUT the conversation "
    "history (i.e. it contains no pronouns or references like 'it', 'that', "
    "'those', 'the previous one' that depend on prior context). "
    "Respond with ONLY 'YES' if it is standalone, or 'NO' if it depends on "
    "the conversation history."
)


def _needs_rewrite(query: str) -> bool:
    """Ask the LLM whether the query is already standalone; skip the (more
    expensive) rewrite call entirely when it is, per the 'rewrite only if
    needed' requirement. Falls back to 'needs rewrite' on any error, which
    is the safer default (a redundant rewrite is harmless; a missing one
    can break retrieval).
    """
    try:
        client = get_groq_client()
        response = client.chat.completions.create(
            model=GROQ_MODEL,
            temperature=0.0,
            messages=[
                {"role": "system", "content": STANDALONE_CHECK_PROMPT},
                {"role": "user", "content": query},
            ],
        )
        verdict = response.choices[0].message.content.strip().upper()
        return not verdict.startswith("YES")
    except Exception:
        logger.exception("Standalone-check failed; defaulting to rewriting")
        return True


def rewrite_query(query: str) -> str:
    """Rewrite the user's query into a standalone question using conversation
    history, but only when the query is not already standalone.
    """
    if not state.memory:
        return query

    if not _needs_rewrite(query):
        logger.info("Query already standalone; skipping rewrite: '%s'", query)
        return query

    client = get_groq_client()
    history_text = "\n".join(
        f"User: {turn['user']}\nAssistant: {turn['assistant']}" for turn in state.memory
    )

    system_prompt = (
        "Given the conversation history and a follow-up question, rewrite the follow-up "
        "question to be a fully standalone question that can be understood without the "
        "conversation history. Preserve the original intent. Respond with ONLY the rewritten "
        "question and nothing else."
    )
    user_prompt = f"Conversation history:\n{history_text}\n\nFollow-up question: {query}\n\nStandalone question:"

    try:
        response = client.chat.completions.create(
            model=GROQ_MODEL,
            temperature=GROQ_TEMPERATURE,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
        )
        rewritten = response.choices[0].message.content.strip()
        logger.info("Rewrote query: '%s' -> '%s'", query, rewritten)
        return rewritten if rewritten else query
    except Exception:
        logger.exception("Query rewriting failed; using original query")
        return query


# ---------------------------------------------------------------------------
# SECTION: Answer Generation
# ---------------------------------------------------------------------------

ANSWER_SYSTEM_PROMPT = (
    "You are DocSensei, an enterprise document assistant. You must answer ONLY using the "
    "provided context extracted from the user's uploaded documents. Never use outside or "
    "general knowledge. If the answer cannot be found in the provided context, respond "
    "exactly with: \"I could not find this information in the uploaded documents.\" "
    "Write the answer naturally, in plain prose. Do NOT include citation markers, source "
    "names, or bracketed references inside the answer text itself; citations are handled "
    "separately."
)


def generate_answer(query: str, context_chunks: list[dict[str, Any]]) -> str:
    """Generate a natural-language answer grounded strictly in the retrieved
    context. Temperature is fixed at GROQ_TEMPERATURE (0.1) to minimize
    hallucination and keep answers deterministic and repeatable.
    """
    client = get_groq_client()

    if not context_chunks:
        return "I could not find this information in the uploaded documents."

    context_text = "\n\n---\n\n".join(c["text"] for c in context_chunks if c["text"].strip())
    if not context_text.strip():
        return "I could not find this information in the uploaded documents."

    user_prompt = f"Context from uploaded documents:\n{context_text}\n\nQuestion: {query}\n\nAnswer:"

    try:
        response = client.chat.completions.create(
            model=GROQ_MODEL,
            temperature=GROQ_TEMPERATURE,
            messages=[
                {"role": "system", "content": ANSWER_SYSTEM_PROMPT},
                {"role": "user", "content": user_prompt},
            ],
        )
        return response.choices[0].message.content.strip()
    except Exception:
        logger.exception("Answer generation failed")
        return "An error occurred while generating the answer. Please try again."


def build_citations(context_chunks: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Build a separate citations list referencing source file, page/slide,
    and chunk id. Citations are never injected into the answer text itself.
    """
    citations = []
    seen = set()
    for chunk in context_chunks:
        key = (chunk.get("source_file"), chunk.get("page"), chunk.get("chunk_id"))
        if key in seen:
            continue
        seen.add(key)
        snippet = chunk.get("text", "")[:220]
        citations.append({
            "source_file": chunk.get("source_file"),
            "page": chunk.get("page"),
            "chunk_id": chunk.get("chunk_id"),
            "file_type": chunk.get("file_type"),
            "snippet": snippet,
        })
    return citations


# ---------------------------------------------------------------------------
# SECTION: Summarization
# ---------------------------------------------------------------------------

SUMMARY_SYSTEM_PROMPT = (
    "You are a precise document summarization assistant. Summarize the given document "
    "into a structured JSON object with EXACTLY these keys:\n"
    '  "executive_summary": 2-4 sentence high-level overview.\n'
    '  "key_topics": array of 3-7 short topic strings.\n'
    '  "important_facts": array of 3-6 concise factual statement strings.\n'
    '  "key_numbers": array of notable figures/statistics/dates as strings '
    "(empty array if the document has none).\n"
    '  "conclusion": 1-3 sentence closing takeaway.\n'
    "Respond with ONLY the JSON object, no preamble, no markdown fences."
)


def _fallback_summary(text: str) -> dict[str, Any]:
    """A minimal structured summary used only if the LLM call fails, so the
    UI never has to special-case a missing summary.
    """
    return {
        "executive_summary": "Summary unavailable due to an internal error.",
        "key_topics": [],
        "important_facts": [],
        "key_numbers": [],
        "conclusion": "",
    }


def summarize_document(filename: str, full_text: str) -> dict[str, Any]:
    """Generate a structured summary for a single uploaded document, covering
    an executive summary, key topics, important facts, key numbers, and a
    conclusion.
    """
    client = get_groq_client()
    truncated_text = full_text[:12000]  # guard against oversized prompts

    user_prompt = f"Document: {filename}\n\nContent:\n{truncated_text}\n\nJSON summary:"

    try:
        response = client.chat.completions.create(
            model=GROQ_MODEL,
            temperature=GROQ_TEMPERATURE,
            messages=[
                {"role": "system", "content": SUMMARY_SYSTEM_PROMPT},
                {"role": "user", "content": user_prompt},
            ],
        )
        raw = response.choices[0].message.content.strip()
        raw = re.sub(r"^```(json)?|```$", "", raw, flags=re.MULTILINE).strip()
        parsed = json.loads(raw)
        return {
            "executive_summary": parsed.get("executive_summary", ""),
            "key_topics": parsed.get("key_topics", []) or [],
            "important_facts": parsed.get("important_facts", []) or [],
            "key_numbers": parsed.get("key_numbers", []) or [],
            "conclusion": parsed.get("conclusion", ""),
        }
    except Exception:
        logger.exception("Summarization failed for file: %s", filename)
        return _fallback_summary(full_text)


# ---------------------------------------------------------------------------
# SECTION: Suggested Questions
# ---------------------------------------------------------------------------

SUGGESTED_QUESTIONS_SYSTEM_PROMPT = (
    "Based on the provided document content, generate exactly 5 insightful questions "
    "that a user could ask and that would be answerable from this content, one for "
    "EACH of the following difficulty/category levels, in this order: "
    "Basic, Intermediate, Advanced, Comparison, Analytical. "
    'Respond ONLY with a JSON array of 5 objects, each shaped as '
    '{"category": "...", "question": "..."}, no preamble, no markdown fences.'
)


def generate_suggested_questions(all_text_sample: str) -> list[dict[str, str]]:
    """Generate exactly 5 suggested questions spanning Basic, Intermediate,
    Advanced, Comparison, and Analytical categories.
    """
    client = get_groq_client()
    truncated_text = all_text_sample[:12000]
    user_prompt = f"Content:\n{truncated_text}"

    try:
        response = client.chat.completions.create(
            model=GROQ_MODEL,
            temperature=GROQ_TEMPERATURE,
            messages=[
                {"role": "system", "content": SUGGESTED_QUESTIONS_SYSTEM_PROMPT},
                {"role": "user", "content": user_prompt},
            ],
        )
        raw = response.choices[0].message.content.strip()
        raw = re.sub(r"^```(json)?|```$", "", raw, flags=re.MULTILINE).strip()
        questions = json.loads(raw)
        if isinstance(questions, list):
            cleaned = []
            for q in questions[:5]:
                if isinstance(q, dict) and q.get("question"):
                    cleaned.append({
                        "category": q.get("category", ""),
                        "question": q["question"],
                    })
                elif isinstance(q, str):
                    cleaned.append({"category": "", "question": q})
            return cleaned
        return []
    except Exception:
        logger.exception("Suggested question generation failed")
        return []


# ---------------------------------------------------------------------------
# SECTION: Helpers
# ---------------------------------------------------------------------------


def allowed_file(filename: str) -> bool:
    ext = os.path.splitext(filename)[1].lower()
    return bool(filename) and ext in ALLOWED_EXTENSIONS


def clear_directory(directory: str) -> None:
    if os.path.isdir(directory):
        for entry in os.listdir(directory):
            path = os.path.join(directory, entry)
            try:
                if os.path.isfile(path) or os.path.islink(path):
                    os.remove(path)
                elif os.path.isdir(path):
                    shutil.rmtree(path)
            except Exception:
                logger.exception("Failed to remove: %s", path)


# ---------------------------------------------------------------------------
# SECTION: Flask Routes
# ---------------------------------------------------------------------------


@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")


@app.route("/upload", methods=["POST"])
def upload():
    try:
        files = request.files.getlist("files")
        if not files:
            return jsonify({"error": "No files provided."}), 400

        valid_files = []
        for f in files:
            if f and allowed_file(f.filename):
                valid_files.append(f)
            else:
                logger.warning("Rejected unsupported file: %s", getattr(f, "filename", None))

        if not valid_files:
            return jsonify({"error": "No supported files (.pdf, .docx, .pptx, .txt) were provided."}), 400

        logger.info("Received %d valid file(s) for upload", len(valid_files))

        # Delete only the PREVIOUS namespace's vectors (never a global delete),
        # then fully reset in-memory session state for the new batch.
        delete_namespace()
        state.reset()
        clear_directory(UPLOAD_FOLDER)
        clear_directory(TEMP_FOLDER)

        namespace = new_namespace()

        saved_paths: dict[str, str] = {}
        for f in valid_files:
            safe_name = f"{uuid.uuid4().hex}_{secure_filename(f.filename)}"
            save_path = os.path.join(UPLOAD_FOLDER, safe_name)
            f.save(save_path)
            saved_paths[f.filename] = save_path

        all_chunks: list[dict[str, Any]] = []
        failed_files: list[str] = []
        for original_name, path in saved_paths.items():
            ext = os.path.splitext(original_name)[1].lower()
            loader = LOADER_MAP.get(ext)
            if not loader:
                failed_files.append(original_name)
                continue

            try:
                pages = loader(path)
                full_text = "\n\n".join(clean_text(p["text"]) for p in pages)
                if not full_text.strip():
                    logger.warning("No extractable text in file: %s", original_name)
                    failed_files.append(original_name)
                    continue

                state.raw_documents[original_name] = full_text
                file_chunks = build_chunks_for_file(original_name, ext.lstrip("."), pages)
                all_chunks.extend(file_chunks)
                state.uploaded_files.append(original_name)
            except Exception:
                # Never let one bad document abort the whole batch.
                logger.exception("Failed to process file '%s'; continuing with remaining files", original_name)
                failed_files.append(original_name)
                continue

        if not all_chunks:
            return jsonify({"error": "No extractable text found in the uploaded files."}), 400

        state.chunks = all_chunks
        state.namespace = namespace
        state.previous_namespace = namespace  # becomes "previous" on the NEXT upload

        # Embed & upsert to the new namespace
        upsert_chunks_to_pinecone(state.chunks, namespace)

        # Build BM25 index
        build_bm25_index()

        # Generate per-document structured summaries (one bad summary doesn't block others)
        summaries: dict[str, Any] = {}
        for filename, text in state.raw_documents.items():
            try:
                summaries[filename] = summarize_document(filename, text)
            except Exception:
                logger.exception("Summary generation failed for '%s'", filename)
                summaries[filename] = _fallback_summary(text)
        state.summaries = summaries

        # Generate suggested questions from a combined sample of all docs
        combined_sample = "\n\n".join(state.raw_documents.values())
        try:
            suggested = generate_suggested_questions(combined_sample)
        except Exception:
            logger.exception("Suggested question generation failed for upload batch")
            suggested = []
        state.suggested_questions = suggested

        response_payload = {
            "status": "success",
            "files": state.uploaded_files,
            "failed_files": failed_files,
            "chunk_count": len(state.chunks),
            "namespace": state.namespace,
            "summaries": state.summaries,
            "suggested_questions": state.suggested_questions,
            "timestamp": datetime.utcnow().isoformat(),
        }
        return jsonify(response_payload)

    except Exception:
        logger.exception("Upload failed")
        return jsonify({"error": "An internal error occurred while processing the upload."}), 500


@app.route("/chat", methods=["POST"])
def chat():
    try:
        payload = request.get_json(force=True, silent=True) or {}
        query = (payload.get("message") or "").strip()

        if not query:
            return jsonify({"error": "No message provided."}), 400

        if not state.chunks:
            return jsonify({
                "answer": "Please upload one or more documents before asking questions.",
                "citations": [],
            })

        # 1. Rewrite query using conversation history (skipped if already standalone)
        standalone_query = rewrite_query(query)

        # 2. Adaptive retrieval: LLM router -> hybrid retrieve -> CrossEncoder rerank
        retrieved_chunks, query_type = adaptive_retrieve(standalone_query)

        # 3. Contextual compression (top-N only)
        compressed_chunks = compress_context(standalone_query, retrieved_chunks)

        # 4. Generate answer
        answer = generate_answer(standalone_query, compressed_chunks)

        # 5. Build citations (returned separately, never injected into the answer)
        citations = build_citations(compressed_chunks)

        # 6. Update conversation memory
        state.memory.append({"user": query, "assistant": answer})

        return jsonify({
            "answer": answer,
            "citations": citations,
            "query_type": query_type,
            "rewritten_query": standalone_query,
        })

    except Exception:
        logger.exception("Chat request failed")
        return jsonify({"error": "An internal error occurred while generating the response."}), 500


@app.route("/clear", methods=["POST"])
def clear():
    """Clear conversation memory and derived artifacts, but keep the current
    namespace's vectors/documents intact (lighter-weight than /reset).
    """
    try:
        state.memory.clear()
        logger.info("Conversation memory cleared")
        return jsonify({"status": "cleared", "timestamp": datetime.utcnow().isoformat()})
    except Exception:
        logger.exception("Clear operation failed")
        return jsonify({"error": "An internal error occurred while clearing state."}), 500


@app.route("/reset", methods=["POST"])
def reset():
    """Completely clear namespace, memory, chunks, BM25, and temporary files."""
    try:
        clear_directory(UPLOAD_FOLDER)
        clear_directory(TEMP_FOLDER)
        delete_namespace()
        delete_namespace()
        state.reset()
        logger.info("Application fully reset")
        return jsonify({"status": "reset", "timestamp": datetime.utcnow().isoformat()})
    except Exception:
        logger.exception("Reset operation failed")
        return jsonify({"error": "An internal error occurred while resetting the application."}), 500


# ---------------------------------------------------------------------------
# Global error handlers
# ---------------------------------------------------------------------------


@app.errorhandler(404)
def not_found(_error):
    return jsonify({"error": "Endpoint not found."}), 404


@app.errorhandler(413)
def payload_too_large(_error):
    return jsonify({"error": "Upload exceeds the maximum allowed size (500 MB)."}), 413


@app.errorhandler(500)
def server_error(_error):
    logger.exception("Unhandled server error")
    return jsonify({"error": "Internal server error."}), 500


@app.errorhandler(Exception)
def handle_uncaught_exception(error):
    logger.exception("Uncaught exception: %s", error)
    return jsonify({"error": "An unexpected error occurred."}), 500


# ---------------------------------------------------------------------------
# Entrypoint
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    debug_mode = os.environ.get("FLASK_DEBUG", "false").lower() == "true"
    logger.info("Starting DocSensei on port %d (debug=%s)", port, debug_mode)
    ocr_is_available()  # log OCR availability once at startup, not on first request
    app.run(host="0.0.0.0", port=port, debug=debug_mode)
